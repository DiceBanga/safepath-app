#!/usr/bin/env python3
"""
Integrate images into existing presentation slides at appropriate positions
"""

from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.enum.text import PP_ALIGN
import os

# Open existing presentation
prs = Presentation('/home/ubuntu/.openclaw/workspace/safepath-proposal/SafePath_Presentation.pptx')

# Image directory
img_dir = '/home/ubuntu/.openclaw/workspace/safepath-proposal/img'

print(f"Total slides: {len(prs.slides)}")

# Slide 1: Title Slide - Add hero image as background/side image
slide1 = prs.slides[0]
try:
    slide1.shapes.add_picture(
        f'{img_dir}/safepath-hero.jpg',
        Inches(6), Inches(1.5), width=Inches(3.5)
    )
    print("✓ Added hero image to title slide")
except Exception as e:
    print(f"✗ Error adding to slide 1: {e}")

# Slide 2: Problem Statement - Add hazard detection concept
slide2 = prs.slides[1]
try:
    slide2.shapes.add_picture(
        f'{img_dir}/hazard-detection.jpg',
        Inches(6), Inches(1.5), width=Inches(3.5)
    )
    print("✓ Added hazard diagram to problem statement slide")
except Exception as e:
    print(f"✗ Error adding to slide 2: {e}")

# Slide 3: Objectives - Small icon or leave as is
print("✓ Objectives slide - left text-only")

# Slide 4: Technical Approach - Add pipeline diagram
slide4 = prs.slides[3]
try:
    slide4.shapes.add_picture(
        f'{img_dir}/pipeline-diagram.jpg',
        Inches(5.5), Inches(2), width=Inches(4.2)
    )
    print("✓ Added pipeline diagram to technical approach slide")
except Exception as e:
    print(f"✗ Error adding to slide 4: {e}")

# Slide 5: Dataset Strategy - Add semantic segmentation illustration
slide5 = prs.slides[4]
try:
    slide5.shapes.add_picture(
        f'{img_dir}/semantic-segmentation.jpg',
        Inches(5.5), Inches(1.8), width=Inches(4)
    )
    print("✓ Added semantic segmentation image to dataset strategy slide")
except Exception as e:
    print(f"✗ Error adding to slide 5: {e}")

# Slide 6: Timeline - Leave as is (text-focused)
print("✓ Timeline slide - left text-only")

# Slide 7: Risk Assessment - Leave as is (table)
print("✓ Risk assessment slide - left as is")

# Slide 8: Success Metrics - Leave as is
print("✓ Success metrics slide - left as is")

# Slide 9: Academic Context - Leave as is
print("✓ Academic context slide - left as is")

# Slide 10: Next Steps - Leave as is
print("✓ Next steps slide - left as is")

# Add cover art options at the end
def add_cover_slide(title, image_path, description):
    """Add a cover art slide"""
    blank_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(blank_layout)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(28)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Image
    slide.shapes.add_picture(image_path, Inches(2.5), Inches(1.2), width=Inches(5))

    # Description
    desc_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(9), Inches(0.4))
    desc_frame = desc_box.text_frame
    desc_frame.text = description
    desc_frame.paragraphs[0].font.size = Pt(14)
    desc_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    return slide

# Add Cover Art Option 1
add_cover_slide(
    "Cover Art Option 1 - Professional Blue",
    f'{img_dir}/cover-art-1.jpg',
    "Modern minimalist design with blue accents"
)

# Add Cover Art Option 2 (Nano Banana Yellow)
add_cover_slide(
    "Cover Art Option 2 - Nano Banana Yellow",
    f'{img_dir}/cover-art-2.jpg',
    "Stylized design with yellow accent theme"
)

# Save updated presentation
prs.save('/home/ubuntu/.openclaw/workspace/safepath-proposal/SafePath_Presentation_Final.pptx')
print("\n" + "="*50)
print("Presentation updated successfully!")
print(f"Saved as: SafePath_Presentation_Final.pptx")
print("="*50)
print(f"\nTotal slides: {len(prs.slides)}")
print("\nImage integration summary:")
print("  • Slide 1 (Title): Hero image added")
print("  • Slide 2 (Problem): Hazard diagram added")
print("  • Slide 3 (Objectives): Text-only (clean)")
print("  • Slide 4 (Technical): Pipeline diagram added")
print("  • Slide 5 (Dataset): Segmentation image added")
print("  • Slides 6-10: Text/Table (kept clean)")
print("  • Bonus: 2 cover art options added at end")
