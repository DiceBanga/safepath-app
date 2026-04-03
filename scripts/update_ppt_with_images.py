#!/usr/bin/env python3
"""
Update PowerPoint presentation with images
"""

from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.enum.text import PP_ALIGN
import os

# Open existing presentation
prs = Presentation('/home/ubuntu/.openclaw/workspace/safepath-proposal/SafePath_Presentation.pptx')

# Image directory
img_dir = '/home/ubuntu/.openclaw/workspace/safepath-proposal/img'

def add_slide_with_image(title, subtitle, image_path, img_width=Inches(8), img_left=Inches(1), img_top=Inches(1.5)):
    """Add a slide with a centered image"""
    blank_layout = prs.slide_layouts[5]  # Blank
    slide = prs.slides.add_slide(blank_layout)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(28)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Add image
    if os.path.exists(image_path):
        slide.shapes.add_picture(image_path, img_left, img_top, width=img_width)

    # Subtitle
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(9), Inches(0.4))
        sub_frame = sub_box.text_frame
        sub_frame.text = subtitle
        sub_frame.paragraphs[0].font.size = Pt(16)
        sub_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    return slide

# Create new image slides
# Image Slide 1: Hero
add_slide_with_image(
    "SafePath - AI Navigation Solution",
    "Real-time hazard detection for safer navigation",
    f'{img_dir}/safepath-hero.jpg',
    img_width=Inches(8.5)
)

# Image Slide 2: Hazard Detection
add_slide_with_image(
    "Hazard Detection Concept",
    "Identifying obstacles in low-light environments",
    f'{img_dir}/hazard-detection.jpg'
)

# Image Slide 3: Pipeline
add_slide_with_image(
    "Processing Pipeline Architecture",
    "From camera input to PDF report generation",
    f'{img_dir}/pipeline-diagram.jpg'
)

# Image Slide 4: Semantic Segmentation
add_slide_with_image(
    "Semantic Segmentation Output",
    "Scene understanding through pixel-level classification",
    f'{img_dir}/semantic-segmentation.jpg'
)

# Cover Art Draft 1
add_slide_with_image(
    "Cover Art Draft - Professional Blue",
    "Modern minimalist design",
    f'{img_dir}/cover-art-1.jpg'
)

# Cover Art Draft 2
add_slide_with_image(
    "Cover Art Draft - Nano Banana Yellow",
    "Stylized with yellow accent",
    f'{img_dir}/cover-art-2.jpg'
)

# Save updated presentation
prs.save('/home/ubuntu/.openclaw/workspace/safepath-proposal/SafePath_Presentation_with_Images.pptx')
print("Updated presentation saved: SafePath_Presentation_with_Images.pptx")
