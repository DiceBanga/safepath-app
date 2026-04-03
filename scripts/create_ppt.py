#!/usr/bin/env python3
"""
Create PowerPoint presentation from SafePath slides
"""

from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.enum.text import PP_ALIGN

# Create presentation
prs = Presentation()

# Set slide dimensions to 16:9
prs.slide_width = Inches(10)
prs.slide_height = Inches(5.625)

# Helper function to create title slide
def create_title_slide(title, subtitle, date):
    slide_layout = prs.slide_layouts[0]  # Title slide
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title

    # Find subtitle placeholder (usually index 1) or add text box
    for shape in slide.placeholders:
        if shape.placeholder_format.type == 2:  # Subtitle type
            shape.text = subtitle
            break

    # Add date manually as a text box
    date_box = slide.shapes.add_textbox(Inches(5), Inches(4.5), Inches(4), Inches(0.5))
    date_frame = date_box.text_frame
    date_frame.text = date
    date_frame.paragraphs[0].font.size = Pt(12)
    date_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT

    return slide

# Helper function to create content slide with title and bullets
def create_content_slide(title, bullets):
    slide_layout = prs.slide_layouts[1]  # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title

    text_box = slide.placeholders[1].text_frame
    text_box.clear()

    for bullet in bullets:
        p = text_box.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.size = Pt(18)

    return slide

# Helper function to create slide with title and two columns
def create_two_column_slide(title, left_col, right_col):
    slide_layout = prs.slide_layouts[5]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Left column
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(4.2), Inches(4))
    left_frame = left_box.text_frame
    for item in left_col:
        p = left_frame.add_paragraph()
        p.text = item
        p.font.size = Pt(18)

    # Right column
    right_box = slide.shapes.add_textbox(Inches(5.3), Inches(1.3), Inches(4.2), Inches(4))
    right_frame = right_box.text_frame
    for item in right_col:
        p = right_frame.add_paragraph()
        p.text = item
        p.font.size = Pt(18)

    return slide

# Helper function to create table slide
def create_table_slide(title, headers, rows_data):
    slide_layout = prs.slide_layouts[5]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Table
    shape = slide.shapes.add_table(rows=len(rows_data) + 1, cols=len(headers),
                                  left=Inches(0.5), top=Inches(1.3),
                                  width=Inches(9), height=Inches(3.5))
    table = shape.table

    # Header row
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(14)
        cell.fill.solid()

    # Data rows
    for row_idx, row_data in enumerate(rows_data, start=1):
        for col_idx, cell_data in enumerate(row_data):
            cell = table.cell(row_idx, col_idx)
            cell.text = str(cell_data)
            cell.text_frame.paragraphs[0].font.size = Pt(12)

    return slide

# Create slides
# Slide 1: Title Slide
create_title_slide(
    "SafePath",
    "AI-Powered Scene Analysis Application",
    "Dice | AI688 - Image and Vision Computing\nDate: February 26, 2026"
)

# Slide 2: Problem Statement
create_content_slide(
    "Problem Statement - Navigating Dark",
    [
        "Safety risks in low-light environments (trips, falls, collisions)",
        "Limited visibility hinders hazard detection",
        "Need for a hands-free, automated visual aid"
    ]
)

# Slide 3: Project Objectives
create_content_slide(
    "Project Objectives - Core Goals",
    [
        "Hazard Detection: Identify obstacles (potholes, holes, poles) with high accuracy",
        "Low-Light Optimization: Tune specifically for dark environments",
        "Mobile Deployment: Efficient processing on Samsung Galaxy Z Fold 6",
        "Reporting: Automated PDF generation for incident logs"
    ]
)

# Slide 4: Technical Approach
# Hardware Target
left_col = [
    "Samsung Galaxy Z Fold 6",
    "Superior low-light sensors",
    "NPU acceleration for efficiency"
]
# Vision Model
right_col = [
    "DeepLabV3+ (Pre-trained on Cityscapes)",
    "Fine-tuning on custom hazard datasets",
    "Semantic Segmentation for path extraction"
]
create_two_column_slide("Technical Approach - Hardware Target & Vision Model", left_col, right_col)

# Software Stack
create_content_slide(
    "Technical Approach - Software Stack",
    [
        "Python 3.9+",
        "PyTorch (Inference)",
        "OpenCV (Image/Video processing)",
        "ReportLab (PDF generation)"
    ]
)

# Slide 5: Dataset Strategy - Training Sources
left_col = [
    "BDD100K (Driving scenes)",
    "Cityscapes (Urban hazards)",
    "Custom Low-Light Set (Nighttime simulation)"
]
right_col = [
    "Brightness reduction",
    "Gaussian noise",
    "Motion blur (to simulate walking)"
]
create_two_column_slide("Dataset Strategy - Training Sources & Augmentation", left_col, right_col)

# Slide 6: Deliverables & Timeline
create_content_slide(
    "Deliverables & Timeline",
    [
        "Phase 1: Planning & Research (Due Feb 27)",
        "  • Proposal Submission",
        "  • Environment Setup",
        "Phase 2: Model Development (Due Mar 15)",
        "  • Model Selection & Fine-tuning",
        "  • Initial Training & Validation",
        "Phase 3: Proof of Concept (Due Apr 2)",
        "  • Mobile Application Development",
        "  • Path Detection & Warning Overlay",
        "  • PDF Report Generation",
        "Phase 4: Finalization (Due May 7)",
        "  • Performance Optimization",
        "  • Final Presentation",
        "  • Documentation"
    ]
)

# Slide 7: Risk Assessment
create_table_slide(
    "Risk Assessment - Potential Risks",
    ["Risk", "Severity", "Mitigation"],
    [
        ["Dataset mismatch", "High", "Early testing & augmentation"],
        ["Low inference speed", "Medium", "Model quantization (FP16) and utilize NPU"],
        ["Poor night detection", "High", "Z Fold 6 specific tuning"]
    ]
)

# Slide 8: Success Metrics
create_content_slide(
    "Success Metrics - Measurable Targets",
    [
        "Detection Accuracy: > 85%",
        "Inference Speed: > 5 FPS",
        "Report Generation: < 5 sec processing time"
    ]
)

# Slide 9: Academic Context
create_content_slide(
    "Academic Context - Course Information",
    [
        "Course: AI688 - Image and Vision Computing",
        "Section: Section 1",
        "Individual Project"
    ]
)

# Slide 10: Next Steps
create_content_slide(
    "Next Steps",
    [
        "[ ] Submit Project Proposal (Today, 6 PM)",
        "[ ] Begin dataset research",
        "[ ] Set up deep learning environment",
        "[ ] Develop initial proof of concept"
    ]
)

# Save presentation
prs.save('/home/ubuntu/.openclaw/workspace/safepath-proposal/SafePath_Presentation.pptx')
print("Presentation created: SafePath_Presentation.pptx")
